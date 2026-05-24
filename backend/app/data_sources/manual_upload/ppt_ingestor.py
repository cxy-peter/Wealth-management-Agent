from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any


def preview_ppt(source: bytes | str | Path, sample_items: int = 40) -> dict[str, Any]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        return {"parser_status": "unsupported_or_optional", "file_type": "ppt", "error_type": exc.__class__.__name__, "texts": []}
    try:
        presentation = Presentation(BytesIO(source) if isinstance(source, bytes) else source)
        texts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text.strip())
                if len(texts) >= sample_items:
                    break
            if len(texts) >= sample_items:
                break
        return {"parser_status": "parsed", "file_type": "ppt", "texts": texts}
    except Exception as exc:
        return {"parser_status": "failed", "file_type": "ppt", "error_type": exc.__class__.__name__, "texts": []}

